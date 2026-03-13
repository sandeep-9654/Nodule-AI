from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullet_points[0] if bullet_points else ""
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
            
    return slide

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Nodule AI"
    subtitle.text = "Volumetric CT Scan Analysis for Lung Nodule Detection\n\nTeam 14"
    
    # 2. Introduction
    add_slide(prs, "Introduction", [
        "Lung cancer is one of the leading causes of cancer-related mortality worldwide.",
        "Early detection of pulmonary nodules significantly increases survival rates.",
        "Computed Tomography (CT) scans are the standard imaging modality for screening.",
        "Manual inspection of volumetric CT scans is time-consuming and prone to human error.",
        "Nodule AI provides an automated, AI-driven solution to assist radiologists in detecting and classifying lung nodules."
    ])
    
    # 3. Problem Statement
    add_slide(prs, "Problem Statement", [
        "Class imbalance in medical datasets often leads to biased AI models.",
        "High computational costs associated with processing full 3D isotropic medical volumes.",
        "Legacy datasets (e.g., LUNA-16) often cause severe hardware exhaustion during local training.",
        "Lack of interactive, explainable, and visually intuitive 3D visualization for clinicians."
    ])
    
    # 4. Existing Systems vs. Limitations
    add_slide(prs, "Existing Systems", [
        "Traditional Computer-Aided Diagnosis (CAD) Systems:",
        " • Rely heavily on hand-crafted features (shape, texture, intensity).",
        " • High false-positive rates and brittle generalization to new data.",
        "Standard Full-Scale 3D CNNs:",
        " • Demand massive GPU memory (VRAM).",
        " • Difficult to deploy in localized or resource-constrained hospital environments.",
        "Manual Review by Radiologists:",
        " • Slice-by-slice 2D inspection of a 3D scan is cognitively taxing."
    ])

    # 5. Proposed Solution
    add_slide(prs, "Proposed Solution: Nodule AI", [
        "Migrated to the pre-extracted, balanced LUNA-22 (ISMI) 3D Patches dataset to solve class imbalance.",
        "Developed a deeply optimized Lightweight U-Net 3D model utilizing Depthwise Separable Convolutions.",
        "Reduced parameter count to ~155K (< 1 MB footprint) while preserving spatial context.",
        "A full-stack web application (Next.js + Flask) providing real-time 3D rendering and clinical reasoning."
    ])

    # 6. Objectives
    add_slide(prs, "System Objectives", [
        "Provide real-time, highly accurate 3D lung nodule classification (Benign vs. Malignant).",
        "Ensure the architecture is lightweight enough to run inference on edge devices/CPUs.",
        "Deliver a visually stunning, CT-authentic 3D interface for radiologists.",
        "Generate rich clinical descriptions detailing nodule size, texture, and recommended clinical actions."
    ])

    # 7. Data Collection and Preprocessing
    add_slide(prs, "Data Collection & Preprocessing", [
        "Dataset: LUNA-22 (ISMI) dataset containing 1,176 pre-extracted 3D patches.",
        "Data Format: Medical '.nii.gz' volumetric arrays.",
        "Preprocessing Steps:",
        " • Voxel Intensity Normalization to a mapped [0, 1] scale.",
        " • Spatial Resampling & Trilinear Interpolation to a standardized (64×64×32) dimension.",
        " • Automated data-discovery during cloud training to flexibly locate dataset artifacts."
    ])

    # 8. Feasibility Study
    add_slide(prs, "Feasibility Study", [
        "Technical Feasibility:",
        " • Built on robust open-source stacks: PyTorch, Flask, Next.js, and Three.js.",
        " • Lightweight models eliminate the need for enterprise-grade hospital GPUs for inference.",
        "Operational Feasibility:",
        " • Web-browser deployment requires zero installation for clinicians.",
        "Economic Feasibility:",
        " • Utilizes free tier cloud environments (Kaggle) for heavy compute (training).",
        " • Local inference is practically free, running efficiently on consumer hardware."
    ])

    # 9. Process Model Used
    add_slide(prs, "Process Model Used", [
        "Incremental & Iterative Development Model:",
        " • Phase 1: Environment Setup & Legacy Code Purge.",
        " • Phase 2: Neural Architecture Redesign (Lightweight U-Net 3D).",
        " • Phase 3: Backend API Refactoring (Raw Sigmoid + Mesh Extraction).",
        " • Phase 4: Cloud Training Pipeline Generation (Combined Loss, AMP).",
        " • Phase 5: UI/UX Aesthetic and 3D Visualizer enhancements."
    ])

    # 10. System Architecture
    add_slide(prs, "System Architecture", [
        "1. Client Layer (Next.js & Tailwind CSS):",
        "   • Handles user `.npy` uploads and highly responsive layout.",
        "2. Visualization Layer (Three.js & React-Three-Fiber):",
        "   • Procedural rendering of CT-style lung shapes and localized glowing nodules.",
        "3. API Gateway (Flask RESTful API):",
        "   • Handles preprocessing, inference requests, and 3D mesh marching-cubes generation.",
        "4. Deep Learning Engine (PyTorch):",
        "   • Executes the rapid forward-pass on the frozen Lightweight U-Net 3D graph."
    ])

    # 11. Development Phases Summary
    add_slide(prs, "Development Phases", [
        "Phase 1 - Purging Legacy Bloat: Removed heavyweight UNet (90MB) and brittle hardcoded loops.",
        "Phase 2 - Architecture Rebuild: Enforced strict Encoder/Decoder block hierarchies.",
        "Phase 3 - Inference Engine Optimization: Removed synthetic predictions and piecewise probability curves.",
        "Phase 4 - Cloud Training Ready: Scripted Kaggle execution loops utilizing mixed precision (AMP) and safe Dice metrics."
    ])

    # 12. Object Detection Model: Lightweight U-Net 3D
    add_slide(prs, "Core Model: Lightweight U-Net 3D", [
        "Why U-Net?",
        " • The U-shape architecture excels in medical imaging by concatenating deep semantic features with shallow spatial features via skip connections.",
        "Why Lightweight?",
        " • Standard 3D convolutions require massive FLOPs.",
        " • Nodule AI utilizes Depthwise Separable 3D Convolutions, cutting computations drastically.",
        " • Final model size is ~0.59 MB (154,812 parameters)."
    ])

    # 13. Training Strategy & Optimization
    add_slide(prs, "Training Strategy", [
        "Loss Function used: Combined Loss (BCEWithLogitsLoss + Soft Dice Loss).",
        " • Ensures high pixel-level accuracy while fighting any remaining class imbalances.",
        "Optimizer: AdamW with Cosine Annealing Learning Rate Scheduler.",
        "Bias Trick Activation:",
        " • Final layer bias is initialized heavily negative (-5.0) to prevent early training explosive gradients.",
        "Hardware Acceleration: Automated Mixed Precision (AMP) allowing faster training on Kaggle GPUs."
    ])

    # 14. 3D Volume Visualization
    add_slide(prs, "Interactive 3D Visualization", [
        "Implemented using Three.js and React-Three-Fiber.",
        "Procedural anatomical generation:",
        " • Calculates tapered asymmetric lobes with mediastinal concavities mimicking real lung geometries.",
        "CT-Scan Aesthetics:",
        " • Opaque, semi-translucent soft-tissue rendering.",
        " • Dynamic nodule positioning and internal emissive lighting directly correlating to malignancy."
    ])

    # 15. Performance Metrics
    add_slide(prs, "Performance Metrics Evaluated", [
        "Accuracy - Overall correct benign vs. malignant classifications.",
        "Precision - Reduction of false positives (unnecessary biopsies).",
        "Recall (Sensitivity) - Ensuring malignant nodules are not missed.",
        "F1-Score - Harmonic balance of Precision and Recall.",
        "Safe Dice Coefficient - Specialized metric tracking the spatial overlap effectiveness without crashing on empty frames."
    ])

    # 16. Advantages of the System
    add_slide(prs, "Advantages of the System", [
        "Ultra-Fast Inference: Processing 3D volumes in milliseconds on consumer CPUs.",
        "Clinician-Friendly UX: The dark-mode interface and 3D volume viewer drastically reduces radiologist eye strain and cognitive load.",
        "Scalability: Stateless API architecture allows effortless Docker containerization and cloud scaling.",
        "Explainable Outputs: Generates verbose clinical descriptions analyzing size, margin texture, and recommending follow-up protocols."
    ])

    # 17. Applications
    add_slide(prs, "Applications", [
        "Triage Systems in overcrowded radiology departments.",
        "Second-reader (Concurrent read) to back up primary physicians.",
        "Medical education tool to visually demonstrate nodule characteristics to students.",
        "Rural Health Deployment where expert oncologists are not immediately available."
    ])

    # 18. Expected Results
    add_slide(prs, "Expected Results", [
        "Significantly lower screening turnaround time (TAT).",
        "Reduction in false negative rates due to fatigue in manual radiologist screening.",
        "Democratized access to state-of-the-art volumetric AI analysis via browser.",
        "Enhanced patient outcomes due to earlier, more accurate detection."
    ])

    # 19. Conclusion
    add_slide(prs, "Conclusion", [
        "Nodule AI bridges the gap between complex 3D deep learning models and practical clinical deployment.",
        "By optimizing a Lightweight U-Net 3D and embedding it in a visually rich, zero-install web application, we have created a robust tool for pulmonary analysis.",
        "The system is scalable, fast, and represents a significant step towards intelligent, automated medical diagnosis."
    ])
    
    # 20. End Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = "Questions?"

    prs.save('/Users/sandeepraghavendra/Desktop/project/Nodule_AI_Presentation.pptx')
    print("Presentation created successfully at /Users/sandeepraghavendra/Desktop/project/Nodule_AI_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
